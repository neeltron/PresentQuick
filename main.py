from nltk.corpus import stopwords
from nltk.cluster.util import cosine_distance
import numpy as np
import networkx as nx
import nltk
from flask import Flask, render_template, request
from pptx import Presentation

p = Presentation()

nltk.download('stopwords')
def read_full_text(text):
  td = text.splitlines()
  full_text = td[0].split(".")
  txt_arr = []
  for txt in full_text:
    # print(txt)
    txt_arr.append(txt.replace("[^a-zA-Z]", " ").split(" "))
  txt_arr.pop()
  return txt_arr

def txt_similarity(sent1, sent2, stopwords=None):
  if stopwords is None:
    stopwords = []
 
  sent1 = [w.lower() for w in sent1]
  sent2 = [w.lower() for w in sent2]
  all_words = list(set(sent1 + sent2))
  vector1 = [0] * len(all_words)
  vector2 = [0] * len(all_words)
  for w in sent1:
    if w in stopwords:
      continue
    vector1[all_words.index(w)] += 1

  for w in sent2:
    if w in stopwords:
      continue
    vector2[all_words.index(w)] += 1
  return 1 - cosine_distance(vector1, vector2)
 
def build_similarity_matrix(txt_arr, stop_words):
  similarity_matrix = np.zeros((len(txt_arr), len(txt_arr)))
  for idx1 in range(len(txt_arr)):
    for idx2 in range(len(txt_arr)):
      if idx1 == idx2:
        continue 
      similarity_matrix[idx1][idx2] = txt_similarity(txt_arr[idx1], txt_arr[idx2], stop_words)

  return similarity_matrix


def generate_summary(text, top_n=5):
  stop_words = stopwords.words('english')
  summarize_text = []
  txt_arr =  read_full_text(text)
  txt_similarity_martix = build_similarity_matrix(txt_arr, stop_words)
  txt_similarity_graph = nx.from_numpy_array(txt_similarity_martix)
  scores = nx.pagerank(txt_similarity_graph, max_iter = 1000000)
  ranked_txt = sorted(((scores[i],s) for i,s in enumerate(txt_arr)), reverse=True)   
  for i in range(top_n):
    summarize_text.append(" ".join(ranked_txt[i][1]))
  # summ =  ". ".join(summarize_text)
  return summarize_text

app = Flask('app', template_folder = "templates", static_folder = "static")

@app.route('/')
def hello_world():
  return render_template('index.html')

@app.route('/try')
def longText():
  return render_template('tryitout.html')

@app.route('/try', methods = ['POST'])
def process():
  title_form = request.form['title']
  presenter = request.form['presenter']
  back = request.form['back']
  intro = request.form['intro']
  lit = request.form['literature']
  meth = request.form['method']
  res = request.form['result']
  conc = request.form['conc']
  ack = request.form['ack']
  cite = request.form['cite']
  end = request.form['end']
  background = generate_summary(back, 2)
  introduction = generate_summary(intro, 2)
  literature = generate_summary(lit, 2)
  method = generate_summary(meth, 2)
  result = generate_summary(res, 2)
  conclusion = generate_summary(conc, 2)
  acknowledgement = generate_summary(ack, 2)
  layout = p.slide_layouts[0]
  slide = p.slides.add_slide(layout)
  title = slide.shapes.title
  subtitle = slide.placeholders[1] 
  title.text = title_form
  subtitle.text = presenter
  layout1 = p.slide_layouts[1]
  slide = p.slides.add_slide(layout1)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[1]
  title_shape.text = 'Background'
  tf = body_shape.text_frame
  tf.text = background[0]
  for i in range(1, len(background)):
    p2 = tf.add_paragraph()
    p2.text = background[i]
    p2.level = 0
  layout1 = p.slide_layouts[1]
  slide = p.slides.add_slide(layout1)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[1]
  title_shape.text = 'Introduction'
  tf = body_shape.text_frame
  tf.text = introduction[0]
  for i in range(1, len(introduction)):
    p2 = tf.add_paragraph()
    p2.text = introduction[i]
    p2.level = 0
  layout1 = p.slide_layouts[1]
  slide = p.slides.add_slide(layout1)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[1]
  title_shape.text = 'Related Work'
  tf = body_shape.text_frame
  tf.text = literature[0]
  for i in range(1, len(literature)):
    p2 = tf.add_paragraph()
    p2.text = literature[i]
    p2.level = 0
  layout1 = p.slide_layouts[1]
  slide = p.slides.add_slide(layout1)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[1]
  title_shape.text = 'Methods'
  tf = body_shape.text_frame
  tf.text = method[0]
  for i in range(1, len(method)):
    p2 = tf.add_paragraph()
    p2.text = method[i]
    p2.level = 0

  layout1 = p.slide_layouts[1]
  slide = p.slides.add_slide(layout1)
  shapes = slide.shapes
  title_shape = shapes.title
  body_shape = shapes.placeholders[1]
  title_shape.text = 'Results'
  tf = body_shape.text_frame
  tf.text = result[0]
  for i in range(1, len(result)):
    p2 = tf.add_paragraph()
    p2.text = result[i]
    p2.level = 0

  p.save("slide1.pptx")
  return render_template('index.html')

app.run(host='0.0.0.0', port=8080)
